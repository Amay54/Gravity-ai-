import os
import uuid
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from backend.reporting.chart_generator import ChartGenerator
from backend.schemas.research import ResearchReport, SharedResearchContext


class PPTXGenerator:
    """
    Generates professional, publication-quality PowerPoint slide decks using python-pptx.
    """

    @staticmethod
    def _get_theme_rgb(theme: str) -> tuple[RGBColor, RGBColor]:
        theme = theme.lower()
        if theme == "dark":
            return RGBColor(56, 189, 248), RGBColor(168, 85, 247)
        elif theme == "minimal":
            return RGBColor(55, 65, 81), RGBColor(107, 114, 128)
        elif theme == "corporate":
            return RGBColor(30, 58, 138), RGBColor(71, 85, 105)
        else:  # professional
            return RGBColor(79, 70, 229), RGBColor(124, 58, 237)

    @classmethod
    def generate(
        cls,
        context: SharedResearchContext,
        report_data: ResearchReport,
        session_id: str,
        version: int,
        theme: str = "Professional",
        user_name: str = "Developer",
    ) -> str:
        ppt_dir = "backend/storage/pptx"
        os.makedirs(ppt_dir, exist_ok=True)
        pptx_path = os.path.join(ppt_dir, f"{session_id}_v{version}.pptx")

        prs = Presentation()
        primary_rgb, secondary_rgb = cls._get_theme_rgb(theme)

        # 1. Slide Layout Templates from Slide Master
        # Layout 0: Title Slide, Layout 1: Title and Content, Layout 5: Title Only, Layout 6: Blank
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]
        title_only_layout = prs.slide_layouts[5]
        prs.slide_layouts[6]

        temp_images = []
        temp_dir = "backend/storage/pptx/tmp"
        os.makedirs(temp_dir, exist_ok=True)

        # Helper to add standard slide with title and footer
        def add_slide(title_text: str):
            slide = prs.slides.add_slide(content_layout)
            title_shape = slide.shapes.title
            title_shape.text = title_text
            # Color header text matching theme
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = primary_rgb
                    run.font.name = "Helvetica"
                    run.font.bold = True
            return slide

        # ==========================================
        # SLIDE 1: COVER
        # ==========================================
        slide_1 = prs.slides.add_slide(title_layout)
        title = slide_1.shapes.title
        subtitle = slide_1.placeholders[1]

        title.text = f"Corporate Research Dossier:\n{report_data.company_profile.name.value}"
        subtitle.text = f"Compiled autonomously via GravityAI Multi-Agent System.\nSession ID: {session_id} | Version: v{version} | Date: {datetime.utcnow().strftime('%Y-%m-%d')}"

        # Set colors
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = primary_rgb
                r.font.name = "Helvetica"

        for p in subtitle.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = secondary_rgb
                r.font.name = "Helvetica"

        # ==========================================
        # SLIDE 2: EXECUTIVE SUMMARY
        # ==========================================
        slide_2 = add_slide("Executive Summary & Quality Metrics")
        tf = slide_2.placeholders[1].text_frame
        tf.text = "Key Findings Overview:"

        p = tf.add_paragraph()
        p.text = f"• Description: {report_data.company_profile.description.value[:250]}..."
        p.level = 1

        p = tf.add_paragraph()
        p.text = (
            f"• Overall Research Quality: {report_data.metadata.research_quality_score * 100:.0f}%"
        )
        p.level = 1

        p = tf.add_paragraph()
        p.text = (
            f"• Overall Confidence rating: {report_data.metadata.overall_confidence * 100:.0f}%"
        )
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Source coverage: {report_data.metadata.official_sources} Official | {report_data.metadata.public_sources} Public Sources"
        p.level = 1

        # ==========================================
        # SLIDE 3: COMPANY PROFILE
        # ==========================================
        slide_3 = add_slide("Company Profile & Sitemap details")
        tf = slide_3.placeholders[1].text_frame
        tf.text = f"Company Overview for {report_data.company_profile.name.value}:"

        p = tf.add_paragraph()
        p.text = f"• HQ Location: {report_data.company_profile.hq_location.value}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Founded Year: {report_data.company_profile.founded_year.value}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Key Leadership: {', '.join(report_data.company_profile.key_leadership.value)}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = (
            f"• Sitemap crawled: {'Yes' if report_data.website_analysis.sitemap_found else 'No'}"
        )
        p.level = 1

        # ==========================================
        # SLIDE 4: BUSINESS MODEL
        # ==========================================
        slide_4 = add_slide("Business Model & Segments")
        tf = slide_4.placeholders[1].text_frame
        tf.text = "Strategic business model parameters:"

        p = tf.add_paragraph()
        p.text = (
            f"• Pricing Model: {report_data.financial_analysis.business_model.pricing_model.value}"
        )
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Revenue Streams: {report_data.financial_analysis.business_model.revenue_streams.value}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Customer Segments: {report_data.financial_analysis.business_model.customer_segments.value}"
        p.level = 1

        # ==========================================
        # SLIDE 5: FINANCIAL HIGHLIGHTS (With Chart)
        # ==========================================
        slide_5 = prs.slides.add_slide(title_only_layout)
        slide_5.shapes.title.text = "Financial Highlights & Revenue Growth"

        # Color header text
        for p in slide_5.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = primary_rgb
                r.font.name = "Helvetica"

        # Textbox left, chart right
        txBox = slide_5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = f"Corporate Valuation: {report_data.financial_analysis.valuation.value}"

        p = tf.add_paragraph()
        p.text = f"Revenue Trends: {report_data.financial_analysis.revenue_trends.value}"
        p.space_before = Pt(12)

        p = tf.add_paragraph()
        p.text = f"Funding: {report_data.financial_analysis.funding_rounds.value}"
        p.space_before = Pt(12)

        if report_data.financial_analysis.revenue_chart_data:
            chart_labels = report_data.financial_analysis.revenue_chart_data.get("labels", [])
            chart_values = report_data.financial_analysis.revenue_chart_data.get("data", [])
            if chart_labels and chart_values:
                _, chart_bytes = ChartGenerator.generate_revenue_chart(
                    chart_labels, chart_values, theme
                )

                temp_filename = f"{uuid.uuid4()}.png"
                temp_path = os.path.join(temp_dir, temp_filename)
                with open(temp_path, "wb") as f:
                    f.write(chart_bytes)
                temp_images.append(temp_path)

                slide_5.shapes.add_picture(temp_path, Inches(5.2), Inches(1.5), width=Inches(4.5))

        # ==========================================
        # SLIDE 6: TECHNOLOGY STACK
        # ==========================================
        slide_6 = add_slide("Technology Stack Profile")
        tf = slide_6.placeholders[1].text_frame
        tf.text = "Identified technologies:"

        p = tf.add_paragraph()
        p.text = f"• Frontend: {report_data.tech_stack.frontend_frameworks.value}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Backend & Servers: {report_data.tech_stack.backend_tech.value}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Cloud Providers: {report_data.tech_stack.cloud_providers.value}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Databases: {report_data.tech_stack.databases.value}"
        p.level = 1

        # ==========================================
        # SLIDE 7: HIRING TRENDS (With Chart)
        # ==========================================
        slide_7 = prs.slides.add_slide(title_only_layout)
        slide_7.shapes.title.text = "Hiring Trajectory & Open Positions"
        for p in slide_7.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = primary_rgb
                r.font.name = "Helvetica"

        txBox = slide_7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = f"Hiring Velocity: {report_data.hiring_trends.hiring_velocity.value}"

        p = tf.add_paragraph()
        p.text = f"Open Roles: {report_data.hiring_trends.open_roles.value}"
        p.space_before = Pt(12)

        p = tf.add_paragraph()
        p.text = f"Growing Areas: {report_data.hiring_trends.top_departments.value}"
        p.space_before = Pt(12)

        if report_data.hiring_trends.hiring_chart_data:
            h_labels = report_data.hiring_trends.hiring_chart_data.get("labels", [])
            h_values = report_data.hiring_trends.hiring_chart_data.get("data", [])
            if h_labels and h_values:
                _, chart_bytes = ChartGenerator.generate_hiring_chart(h_labels, h_values, theme)

                temp_filename = f"{uuid.uuid4()}.png"
                temp_path = os.path.join(temp_dir, temp_filename)
                with open(temp_path, "wb") as f:
                    f.write(chart_bytes)
                temp_images.append(temp_path)

                slide_7.shapes.add_picture(temp_path, Inches(5.2), Inches(1.5), width=Inches(4.5))

        # ==========================================
        # SLIDE 8: COMPETITIVE LANDSCAPE (With Table)
        # ==========================================
        slide_8 = prs.slides.add_slide(title_only_layout)
        slide_8.shapes.title.text = "Competitor Landscape Matrix"
        for p in slide_8.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = primary_rgb

        txBox = slide_8.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(1.5))
        txBox.text_frame.text = (
            f"Positioning: {report_data.competitor_analysis.market_positioning.value}"
        )
        txBox.text_frame.word_wrap = True

        # Add Competitor table
        rows = len(report_data.competitor_analysis.direct_competitors) + 1
        table_shape = slide_8.shapes.add_table(
            rows, 3, Inches(0.5), Inches(2.2), Inches(9.0), Inches(3.0)
        )
        table = table_shape.table

        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(4.5)

        table.cell(0, 0).text = "Competitor"
        table.cell(0, 1).text = "Market Share"
        table.cell(0, 2).text = "Advantages"

        for idx, comp in enumerate(report_data.competitor_analysis.direct_competitors):
            name = comp.name if hasattr(comp, "name") else comp.get("name", "N/A")
            focus = comp.focus if hasattr(comp, "focus") else comp.get("focus", "N/A")
            comparison = (
                comp.comparison if hasattr(comp, "comparison") else comp.get("comparison", "N/A")
            )
            table.cell(idx + 1, 0).text = name
            table.cell(idx + 1, 1).text = focus
            table.cell(idx + 1, 2).text = comparison

        # ==========================================
        # SLIDE 9: PATENT ACTIVITY
        # ==========================================
        slide_9 = prs.slides.add_slide(title_only_layout)
        slide_9.shapes.title.text = "Patent Filings & Innovation"
        for p in slide_9.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = primary_rgb

        txBox = slide_9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = f"Patent Count: {report_data.patent_activity.patent_counts.value}"

        p = tf.add_paragraph()
        p.text = f"Filing Trend: {report_data.patent_activity.filing_trends.value}"
        p.space_before = Pt(12)

        p = tf.add_paragraph()
        p.text = f"Themes: {report_data.patent_activity.innovation_themes.value}"
        p.space_before = Pt(12)

        if report_data.patent_activity.patent_chart_data:
            p_labels = report_data.patent_activity.patent_chart_data.get("labels", [])
            p_values = report_data.patent_activity.patent_chart_data.get("data", [])
            if p_labels and p_values:
                _, chart_bytes = ChartGenerator.generate_patent_chart(p_labels, p_values, theme)

                temp_filename = f"{uuid.uuid4()}.png"
                temp_path = os.path.join(temp_dir, temp_filename)
                with open(temp_path, "wb") as f:
                    f.write(chart_bytes)
                temp_images.append(temp_path)

                slide_9.shapes.add_picture(temp_path, Inches(5.2), Inches(1.5), width=Inches(4.5))

        # ==========================================
        # SLIDE 10: SWOT ANALYSIS
        # ==========================================
        slide_10 = add_slide("SWOT Matrix Evaluation")
        tf = slide_10.placeholders[1].text_frame
        tf.text = "SWOT Highlights:"

        swot = report_data.swot_matrix
        p = tf.add_paragraph()
        p.text = f"• Strengths: {', '.join(swot.strengths[:3])}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Weaknesses: {', '.join(swot.weaknesses[:3])}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Opportunities: {', '.join(swot.opportunities[:3])}"
        p.level = 1

        p = tf.add_paragraph()
        p.text = f"• Threats: {', '.join(swot.threats[:3])}"
        p.level = 1

        # ==========================================
        # SLIDE 11: STRATEGIC RECOMMENDATIONS
        # ==========================================
        slide_11 = add_slide("Strategic Recommendations")
        tf = slide_11.placeholders[1].text_frame
        tf.text = "Actionable corporate guidance items:"

        for idx, rec in enumerate(report_data.strategic_recommendations[:4]):
            p = tf.add_paragraph()
            p.text = f"• Recommendation {idx + 1}: {rec}"
            p.level = 1

        # ==========================================
        # SLIDE 12: CLOSING
        # ==========================================
        slide_12 = prs.slides.add_slide(title_layout)
        title = slide_12.shapes.title
        subtitle = slide_12.placeholders[1]

        title.text = "Thank You"
        subtitle.text = "GravityAI Enterprise Research Operating System\nFor inquiries contact research@gravityai.com"

        for p in title.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = primary_rgb
                r.font.name = "Helvetica"

        for p in subtitle.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = secondary_rgb
                r.font.name = "Helvetica"

        prs.save(pptx_path)

        # Clean up temp images
        for temp_img in temp_images:
            try:
                os.remove(temp_img)
            except Exception:
                pass

        return pptx_path
